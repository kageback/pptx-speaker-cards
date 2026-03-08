#!/usr/bin/env python3
"""
Speaker Cards Generator
Extracts speaker notes from PowerPoint presentations and generates PDF speaker cards.
"""

import argparse
import sys
import os
import tempfile
from pathlib import Path
from dataclasses import dataclass
from typing import List, Optional, Tuple

import requests

from pptx import Presentation
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.units import mm
from reportlab.lib.colors import grey
from reportlab.pdfgen import canvas
from reportlab.platypus import Paragraph
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.enums import TA_LEFT

try:
    from . import __version__
except ImportError:
    __version__ = "1.0.0"


@dataclass
class FormattedRun:
    """Text run with formatting."""
    text: str
    bold: bool = False
    italic: bool = False


@dataclass
class FormattedParagraph:
    """Paragraph with runs and bullet info."""
    runs: List[FormattedRun]
    level: int = 0  # Bullet/indent level
    has_bullet: bool = False  # Whether this paragraph has a bullet point

    def to_html(self) -> str:
        """Convert to HTML markup for ReportLab."""
        html_parts = []

        # Add bullet if present
        if self.has_bullet:
            indent = '&nbsp;' * (self.level * 4)
            html_parts.append(indent + '• ')

        # Process each run
        for run in self.runs:
            text = escape_html(run.text)

            # Apply formatting tags
            if run.bold and run.italic:
                text = f'<b><i>{text}</i></b>'
            elif run.bold:
                text = f'<b>{text}</b>'
            elif run.italic:
                text = f'<i>{text}</i>'

            html_parts.append(text)

        return ''.join(html_parts)


@dataclass
class SlideData:
    """Extracted slide data."""
    number: int
    paragraphs: List[FormattedParagraph]
    is_hidden: bool = False


@dataclass
class Card:
    """Represents a single speaker card."""
    title: str  # Always from first line of notes
    paragraphs: List[FormattedParagraph]
    slide_number: str  # Can be "5" or "5+1", "5+2", etc.

    def get_content_html(self) -> str:
        """Get HTML content for the card."""
        parts = []
        for para in self.paragraphs:
            html = para.to_html()
            # All paragraphs treated equally - one line break each
            parts.append(html if html.strip() else '')
            parts.append('<br/>')
        return ''.join(parts)


@dataclass
class Config:
    """Configuration for speaker card generation."""
    # Typography
    title_font_size: float = 12.0
    body_font_size: Optional[float] = None
    min_font_size: float = 6.0

    # Layout
    margin_top: float = 5 * mm
    margin_bottom: float = 5 * mm
    margin_left: float = 5 * mm
    margin_right: float = 5 * mm
    card_padding: float = 3 * mm

    # Features
    show_slide_numbers: bool = True
    include_hidden: bool = False

    # PDF
    output_path: Optional[str] = None


def escape_html(text: str) -> str:
    """Escape special HTML characters and handle tabs."""
    return (text
            .replace('&', '&amp;')
            .replace('<', '&lt;')
            .replace('>', '&gt;')
            .replace('\t', '&nbsp;&nbsp;&nbsp;&nbsp;'))  # Tab = 4 non-breaking spaces


def is_slide_hidden(slide) -> bool:
    """Check if slide is hidden."""
    try:
        show_attr = slide._element.get('show')
        return show_attr == '0'
    except (AttributeError, TypeError):
        return False


def extract_notes_with_formatting(slide) -> List[FormattedParagraph]:
    """Extract notes from slide with formatting preserved."""
    paragraphs = []

    try:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame

        if not text_frame:
            return paragraphs

        for para in text_frame.paragraphs:
            # Extract runs with formatting
            runs = []
            for run in para.runs:
                formatted_run = FormattedRun(
                    text=run.text,
                    bold=run.font.bold is True,
                    italic=run.font.italic is True
                )
                runs.append(formatted_run)

            # Check if paragraph has text content
            has_text = runs and any(r.text.strip() for r in runs)

            # Check if paragraph has a bullet
            has_bullet = False
            try:
                # A paragraph has a bullet if its text has been formatted with bullets
                # We check the _element to see if bullet formatting is applied
                if hasattr(para, '_element'):
                    pPr = para._element.pPr
                    if pPr is not None:
                        buFont = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
                        buChar = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                        buAutoNum = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
                        # If any bullet element exists and buNone is not present, it has a bullet
                        buNone = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                        if (buFont is not None or buChar is not None or buAutoNum is not None) and buNone is None:
                            has_bullet = True
            except:
                pass

            # Add paragraph if it has text, or if it's empty (to preserve spacing)
            if has_text:
                formatted_para = FormattedParagraph(
                    runs=runs,
                    level=para.level if hasattr(para, 'level') else 0,
                    has_bullet=has_bullet
                )
                paragraphs.append(formatted_para)
            elif not paragraphs or paragraphs[-1].runs:  # Add empty para for spacing
                # Create empty paragraph
                formatted_para = FormattedParagraph(
                    runs=[FormattedRun(text='')],
                    level=0,
                    has_bullet=False
                )
                paragraphs.append(formatted_para)

    except Exception as e:
        print(f"Warning: Failed to extract notes: {e}")

    return paragraphs


def extract_slides_from_pptx(pptx_path: str, config: Config) -> List[SlideData]:
    """Extract all slides from PowerPoint presentation."""
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error: Failed to load presentation: {e}")
        sys.exit(1)

    slides = []
    total_slides = len(prs.slides)

    for idx, slide in enumerate(prs.slides, 1):
        print(f"Processing slide {idx}/{total_slides}...")

        # Check if hidden
        is_hidden = is_slide_hidden(slide)
        if is_hidden and not config.include_hidden:
            print(f"  Skipping hidden slide {idx}")
            continue

        # Extract notes
        paragraphs = extract_notes_with_formatting(slide)

        slide_data = SlideData(
            number=idx,
            paragraphs=paragraphs,
            is_hidden=is_hidden
        )

        slides.append(slide_data)

    return slides


def extract_title_from_notes(slide_data: SlideData) -> Tuple[str, List[FormattedParagraph]]:
    """
    Extract title from first line of notes.
    Returns (title, remaining_paragraphs).
    """
    if not slide_data.paragraphs:
        # No notes at all
        return "", []

    # Get first paragraph
    first_para = slide_data.paragraphs[0]
    first_line_text = ''.join(run.text for run in first_para.runs)

    # Split by newline to get first line as title
    lines = first_line_text.split('\n', 1)
    title = lines[0].strip()

    # Remove first line from paragraphs
    if len(lines) > 1:
        # Keep remaining text in first paragraph with same formatting
        remaining_text = lines[1]
        new_runs = [FormattedRun(
            text=remaining_text,
            bold=first_para.runs[0].bold if first_para.runs else False,
            italic=first_para.runs[0].italic if first_para.runs else False
        )]
        new_para = FormattedParagraph(
            runs=new_runs,
            level=first_para.level,
            has_bullet=first_para.has_bullet
        )
        remaining_paras = [new_para] + slide_data.paragraphs[1:]
    else:
        # First line was the entire paragraph, skip it
        remaining_paras = slide_data.paragraphs[1:]

    return title, remaining_paras


def create_paragraph_style(font_size: float) -> ParagraphStyle:
    """Create paragraph style for given font size."""
    return ParagraphStyle(
        'CardBody',
        fontName='Helvetica',
        fontSize=font_size,
        leading=font_size * 1.05,  # Tighter line spacing (was 1.2)
        leftIndent=0,
        rightIndent=0,
        alignment=TA_LEFT,
    )


def measure_text_height(paragraphs: List[FormattedParagraph], font_size: float,
                        content_width: float) -> float:
    """Measure how much vertical space the text needs."""
    if not paragraphs:
        return 0

    style = create_paragraph_style(font_size)
    html_parts = []
    for para in paragraphs:
        html_parts.append(para.to_html())
        html_parts.append('<br/>')

    html_text = ''.join(html_parts)
    paragraph = Paragraph(html_text, style)

    _, height = paragraph.wrap(content_width, 10000)  # Use large height to get full size
    return height


def fit_text_to_card(paragraphs: List[FormattedParagraph], card_width: float,
                     card_height: float, config: Config,
                     has_title: bool) -> Tuple[float, bool]:
    """
    Find font size that fits text in card.
    Returns (font_size, fits_completely).
    """
    # Calculate available content area
    content_width = card_width - 2 * config.card_padding
    title_height = 5 * mm if has_title else 5 * mm  # Minimal spacing - just title height
    content_height = card_height - 2 * config.card_padding - title_height

    # Start with body font size or default
    initial_size = config.body_font_size if config.body_font_size else 10.0
    current_size = initial_size

    while current_size >= config.min_font_size:
        height = measure_text_height(paragraphs, current_size, content_width)

        if height <= content_height:
            if current_size < initial_size:
                print(f"  Warning: Reduced font size from {initial_size:.1f}pt to {current_size:.1f}pt")
            return current_size, True

        # Reduce by 0.5pt
        current_size -= 0.5

    # Still doesn't fit at minimum size
    print(f"  Warning: Text doesn't fit even at minimum font size {config.min_font_size}pt")
    return config.min_font_size, False


def split_paragraphs_for_continuation(paragraphs: List[FormattedParagraph],
                                      card_width: float, card_height: float,
                                      config: Config, font_size: float,
                                      has_title: bool) -> List[List[FormattedParagraph]]:
    """Split paragraphs into groups that fit on separate cards."""
    content_width = card_width - 2 * config.card_padding
    title_height = 5 * mm if has_title else 5 * mm  # Minimal spacing - just title height
    content_height = card_height - 2 * config.card_padding - title_height

    cards = []
    current_card = []
    current_height = 0

    for para in paragraphs:
        # Measure this paragraph
        para_height = measure_text_height([para], font_size, content_width)

        # Check if adding this paragraph would overflow
        if current_height + para_height > content_height and current_card:
            # Start new card
            cards.append(current_card)
            current_card = [para]
            current_height = para_height
        else:
            current_card.append(para)
            current_height += para_height

    # Add final card
    if current_card:
        cards.append(current_card)

    return cards if cards else [[]]


def generate_cards_for_slide(slide_data: SlideData, config: Config,
                             card_width: float, card_height: float) -> List[Card]:
    """Generate one or more cards for a slide."""
    # Extract title from first line of notes
    title, paragraphs = extract_title_from_notes(slide_data)

    # If no content, create empty card with empty title
    if not paragraphs and not title:
        slide_num = str(slide_data.number) if config.show_slide_numbers else ""
        card = Card(
            title="",
            paragraphs=[],
            slide_number=slide_num
        )
        return [card]

    # Try to fit text
    font_size, fits = fit_text_to_card(paragraphs, card_width, card_height, config, has_title=True)

    if fits:
        # Single card
        slide_num = str(slide_data.number) if config.show_slide_numbers else ""
        card = Card(
            title=title,
            paragraphs=paragraphs,
            slide_number=slide_num
        )
        return [card]

    # Need continuation cards
    para_groups = split_paragraphs_for_continuation(paragraphs, card_width, card_height,
                                                     config, font_size, has_title=True)

    cards = []
    for idx, para_group in enumerate(para_groups, 1):
        # Build title with continuation indicator
        if len(para_groups) == 1:
            card_title = title
        else:
            card_title = f"..continued[{title}]" if idx > 1 else title

        # Build slide number with continuation
        if config.show_slide_numbers:
            if len(para_groups) == 1:
                slide_num = str(slide_data.number)
            else:
                slide_num = f"{slide_data.number}+{idx}"
        else:
            slide_num = ""

        card = Card(
            title=card_title,
            paragraphs=para_group,
            slide_number=slide_num
        )
        cards.append(card)

    return cards


def draw_cut_lines(c: canvas.Canvas):
    """Draw dotted cut lines between cards on landscape page."""
    page_size = landscape(A4)
    c.saveState()

    # Set line style
    c.setStrokeColor(grey)
    c.setLineWidth(0.5)
    c.setDash([2, 2], 0)

    # Vertical center line
    x_center = page_size[0] / 2
    c.line(x_center, 0, x_center, page_size[1])

    # Horizontal center line
    y_center = page_size[1] / 2
    c.line(0, y_center, page_size[0], y_center)

    c.restoreState()


def get_card_position(card_index: int) -> Tuple[float, float, float, float]:
    """
    Get position and dimensions for card at given index (0-3) on landscape page.
    Returns (x, y, width, height) where x, y is bottom-left corner.
    """
    page_size = landscape(A4)
    # Card dimensions on landscape page
    card_width = page_size[0] / 2
    card_height = page_size[1] / 2

    # Position in grid
    col = card_index % 2
    row = card_index // 2

    x = col * card_width
    # ReportLab origin is bottom-left, so flip row
    y = page_size[1] - (row + 1) * card_height

    return x, y, card_width, card_height


def render_card(c: canvas.Canvas, card: Card, x: float, y: float,
                width: float, height: float, config: Config, font_size: float):
    """Render a single card at given position (no rotation needed for landscape page)."""
    padding = config.card_padding

    # Title
    if card.title:
        title_font_size = config.title_font_size
        c.setFont('Helvetica-Bold', title_font_size)
        # Title at top of card
        title_y = y + height - padding - title_font_size
        c.drawString(x + padding, title_y, card.title)

    # Slide number
    if card.slide_number:
        number_font_size = 8
        c.setFont('Helvetica', number_font_size)
        # Number at top-right corner
        number_y = y + height - padding - number_font_size
        number_x = x + width - padding - c.stringWidth(card.slide_number, 'Helvetica', number_font_size)
        c.drawString(number_x, number_y, card.slide_number)

    # Content
    if card.paragraphs:
        content_width = width - 2 * padding
        content_x = x + padding
        content_y_top = y + height - padding - (5 * mm if card.title else 5 * mm)

        # Create paragraph style
        style = create_paragraph_style(font_size)

        # Convert to HTML
        html_text = card.get_content_html()

        # Create paragraph
        paragraph = Paragraph(html_text, style)

        # Wrap and draw
        w, h = paragraph.wrap(content_width, height)
        paragraph.drawOn(c, content_x, content_y_top - h)


def render_pdf(cards: List[Card], output_path: str, config: Config):
    """Render all cards to PDF in landscape orientation."""
    page_size = landscape(A4)
    c = canvas.Canvas(output_path, pagesize=page_size)

    # Calculate card dimensions (landscape page, no rotation needed)
    card_width = page_size[0] / 2
    card_height = page_size[1] / 2

    # Determine font size for each card
    card_font_sizes = []
    for card in cards:
        if card.paragraphs:
            has_title = bool(card.title)
            font_size, _ = fit_text_to_card(card.paragraphs, card_width, card_height,
                                           config, has_title)
        else:
            font_size = config.body_font_size if config.body_font_size else 10.0
        card_font_sizes.append(font_size)

    # Render cards (4 per page)
    for page_idx in range(0, len(cards), 4):
        # Draw cut lines
        draw_cut_lines(c)

        # Render up to 4 cards on this page
        for card_idx in range(4):
            global_card_idx = page_idx + card_idx
            if global_card_idx >= len(cards):
                break

            card = cards[global_card_idx]
            font_size = card_font_sizes[global_card_idx]
            x, y, w, h = get_card_position(card_idx)
            render_card(c, card, x, y, w, h, config, font_size)

        # New page if more cards
        if page_idx + 4 < len(cards):
            c.showPage()

    c.save()


def parse_arguments():
    """Parse command-line arguments."""
    parser = argparse.ArgumentParser(
        description='Generate printable speaker cards from PowerPoint presentation notes.\n\n'
                    'Supports both local files and OneDrive/SharePoint public links.\n'
                    'Creates A4 landscape PDF with 4 cards per page (2x2 grid).',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
Examples:
  %(prog)s presentation.pptx
  %(prog)s "https://sharepoint.com/.../file.pptx?e=..."
  %(prog)s slides.pptx --output my-cards.pdf
  %(prog)s slides.pptx --margin-top 10 --margin-left 8
        '''
    )

    parser.add_argument('input',
                       help='PowerPoint file path or OneDrive/SharePoint URL')

    parser.add_argument('--version',
                       action='version',
                       version=f'%(prog)s {__version__}')

    parser.add_argument('--slide-number',
                       choices=['yes', 'no'],
                       default='yes',
                       help='Show slide numbers (default: yes)')

    parser.add_argument('--title-font-size',
                       type=float,
                       default=12.0,
                       metavar='SIZE',
                       help='Title font size in points (default: 12.0)')

    parser.add_argument('--body-font-size',
                       type=float,
                       metavar='SIZE',
                       help='Override body font size (default: auto-fit)')

    parser.add_argument('--include-hidden',
                       action='store_true',
                       help='Include hidden slides')

    parser.add_argument('--output',
                       metavar='PATH',
                       help='Output PDF path (default: [input]_speaker_notes.pdf)')

    parser.add_argument('--margin-top',
                       type=float,
                       default=5,
                       metavar='MM',
                       help='Top margin in mm (default: 5)')

    parser.add_argument('--margin-bottom',
                       type=float,
                       default=5,
                       metavar='MM',
                       help='Bottom margin in mm (default: 5)')

    parser.add_argument('--margin-left',
                       type=float,
                       default=5,
                       metavar='MM',
                       help='Left margin in mm (default: 5)')

    parser.add_argument('--margin-right',
                       type=float,
                       default=5,
                       metavar='MM',
                       help='Right margin in mm (default: 5)')

    return parser.parse_args()


def is_url(input_str: str) -> bool:
    """Check if input string is a URL."""
    return input_str.startswith(('http://', 'https://'))


def download_from_url(url: str) -> str:
    """
    Download PowerPoint file from URL to temporary location.
    Returns path to temporary file.
    """
    print(f"Downloading from URL...")

    # Add download parameter for SharePoint/OneDrive links
    if 'sharepoint.com' in url or 'onedrive.live.com' in url or '1drv.ms' in url:
        separator = '&' if '?' in url else '?'
        url = url + separator + 'download=1'

    # Download with proper headers
    headers = {
        'User-Agent': 'Mozilla/5.0 (compatible; SpeakerCards/1.0)'
    }

    try:
        response = requests.get(url, headers=headers, allow_redirects=True, timeout=60)
        response.raise_for_status()
    except requests.exceptions.RequestException as e:
        print(f"Error: Failed to download from URL: {e}")
        sys.exit(1)

    # Create temporary file
    temp_file = tempfile.NamedTemporaryFile(
        delete=False,
        suffix='.pptx',
        prefix='speaker_cards_'
    )
    temp_file.write(response.content)
    temp_file.close()

    print(f"Downloaded to temporary file: {temp_file.name}")
    return temp_file.name


def main():
    """Main entry point."""
    args = parse_arguments()

    # Detect if input is URL or file path
    if is_url(args.input):
        # Download from URL to temporary file
        pptx_path = download_from_url(args.input)
        cleanup_needed = True
        # Use original filename for output if not specified
        if args.output:
            output_path = args.output
        else:
            # Use current directory for output when input is URL
            output_path = "speaker_notes.pdf"
    else:
        # Validate local file
        input_path = Path(args.input)
        if not input_path.exists():
            print(f"Error: File not found: {args.input}")
            sys.exit(1)

        pptx_path = args.input
        cleanup_needed = False

        # Determine output path for local file
        if args.output:
            output_path = args.output
        else:
            output_path = str(input_path.parent / f"{input_path.stem}_speaker_notes.pdf")

    try:
        # Build config
        config = Config(
            title_font_size=args.title_font_size,
            body_font_size=args.body_font_size,
            margin_top=args.margin_top * mm,
            margin_bottom=args.margin_bottom * mm,
            margin_left=args.margin_left * mm,
            margin_right=args.margin_right * mm,
            show_slide_numbers=(args.slide_number == 'yes'),
            include_hidden=args.include_hidden,
            output_path=output_path
        )

        print(f"Extracting slides from PowerPoint...")
        slides = extract_slides_from_pptx(pptx_path, config)

        if not slides:
            print("Error: No slides found (or all skipped)")
            sys.exit(1)

        print(f"\nGenerating cards for {len(slides)} slides...")
        all_cards = []
        card_width = A4[0] / 2
        card_height = A4[1] / 2

        for slide in slides:
            cards = generate_cards_for_slide(slide, config, card_width, card_height)
            all_cards.extend(cards)

        print(f"\nRendering {len(all_cards)} cards to PDF...")
        render_pdf(all_cards, output_path, config)

        print(f"\n✓ Generated {len(all_cards)} cards from {len(slides)} slides")
        print(f"✓ Output: {output_path}")

    finally:
        # Clean up temporary file if downloaded from URL
        if cleanup_needed and os.path.exists(pptx_path):
            os.unlink(pptx_path)
            print("✓ Cleaned up temporary file")


if __name__ == '__main__':
    main()
